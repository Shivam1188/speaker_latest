import asyncio
from langchain_core.prompts import PromptTemplate
from langchain_ollama import ChatOllama
import aiohttp
from bs4 import BeautifulSoup
from langchain_core.messages import AIMessage, SystemMessage, HumanMessage
from langchain_core.output_parsers import StrOutputParser
from collections import Counter
import ujson as json
import redis
from datetime import datetime, timedelta
from fastapi import FastAPI, Depends, HTTPException, status, Request, WebSocket, WebSocketDisconnect, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from pydub import AudioSegment
from dotenv import load_dotenv
from urllib.parse import parse_qs, unquote
from concurrent.futures import ThreadPoolExecutor
from fastapi.responses import FileResponse, JSONResponse
import asyncio
import time
from pdf2image import convert_from_path
from typing import List
import os, shutil, uuid
from PIL import Image
import google.generativeai as genai
from pinecone import Pinecone
from auth import hash_password, verify_password, create_access_token
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
import torch
import logging
from urllib.parse import parse_qs,unquote
from jose import JWTError, jwt
from ai_speech_module import Topic
from schemas import UserCreate, UserOut, UserUpdate, Token, LoginRequest, ForgotPasswordRequest, GeminiRequest, ChatRequest
from langchain_ollama import OllamaLLM
import re
from typing import List
from fastapi import FastAPI, UploadFile, File, Form
from pdf2image import convert_from_path
from PIL import Image
from docx2pdf import convert as docx_to_pdf
import os, uuid, shutil, logging, tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from fastapi import HTTPException, status
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
# from langchain_community.vectorstores import Pinecone
from PIL import Image, ImageDraw, ImageFont
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
from fastapi import UploadFile, File, Form, HTTPException
import pandas as pd
import docx2txt
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from firestore_models import FirestoreEssay
from firebase import db
from concurrent.futures import ProcessPoolExecutor
process_pool = ProcessPoolExecutor()
import requests
from fastapi import WebSocket, WebSocketDisconnect
from datetime import datetime
import aiofiles
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks
import aiohttp
import tempfile
import asyncio
import logging
from datetime import datetime
from urllib.parse import parse_qs

from fastapi import WebSocket, WebSocketDisconnect
from pydub import AudioSegment
import aiohttp
CPU_API_BASE = "http://13.200.201.10:8000"

scraping_api_key = os.getenv("SCRAPINGDOG_API_KEY")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = os.getenv("PINECONE_ENVIRONMENT")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
EMBEDDING_MODEL_NAME = "embaas/sentence-transformers-e5-large-v2"

SECRET_KEY = "jwt_secret_key"
ALGORITHM="HS256"


logging.basicConfig(
    filename='rag_log.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

load_dotenv()
app = FastAPI(title="FastAPI‑Firebase")

redis_client = redis.StrictRedis(
    host=os.getenv("REDIS_HOST"),
    port=int(os.getenv("REDIS_PORT")),
    username=os.getenv("REDIS_USERNAME"),
    password=os.getenv("REDIS_PASSWORD"),
    decode_responses=True
)

origins = ["https://llm.edusmartai.com","http://localhost:3000","http://localhost:5173"]

app.add_middleware(CORSMiddleware, allow_origins=origins, allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

def get_user_from_redis_session(request: Request):
    token = request.headers.get("Authorization")
    if not token or not token.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token")
    token = token.split(" ")[1]
    session_data = redis_client.get(f"session:{token}")
    if not session_data:
        raise HTTPException(status_code=401, detail="Session expired or invalid")
    return json.loads(session_data)

@app.post("/register", response_model=UserOut)
def register(user: UserCreate):
    user_ref = db.collection("users").where("username", "==", user.username).stream()
    if any(user_ref):
        raise HTTPException(400, "Username or email already exists")

    doc_ref = db.collection("users").add({
        "username": user.username,
        "email": user.email,
        "password": hash_password(user.password)
    })
    user_id = doc_ref[1].id
    return UserOut(id=user_id, username=user.username, email=user.email)

@app.post("/login", response_model=Token)
def login(data: LoginRequest):
    docs = db.collection("users").where("username", "==", data.username).stream()
    user_doc = next(docs, None)
    if not user_doc or not verify_password(data.password, user_doc.to_dict()["password"]):
        raise HTTPException(status.HTTP_401_UNAUTHORIZED, detail="Bad credentials")

    token = create_access_token({"sub": user_doc.id})
    redis_client.setex(f"session:{token}", timedelta(hours=1), json.dumps({"user_id": user_doc.id, "username": data.username}))
    return Token(access_token=token, username=data.username)

@app.get("/logout")
def logout(request: Request):
    token = request.headers.get("Authorization")
    if token and token.startswith("Bearer "):
        redis_client.delete(f"session:{token.split(' ')[1]}")
    return {"detail": "Logged out"}

@app.get("/me", response_model=UserOut)
def me(user=Depends(get_user_from_redis_session)):
    doc = db.collection("users").document(user["user_id"]).get()
    if not doc.exists:
        raise HTTPException(404, "User not found")
    data = doc.to_dict()
    return UserOut(id=doc.id, username=data["username"], email=data["email"])



@app.post("/generate-prompt")
async def generate_prompt(data: GeminiRequest, user=Depends(get_user_from_redis_session)):
    url = f"https://en.wikipedia.org/wiki/{data.topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()
                    soup = BeautifulSoup(html, "html.parser")
                    for script in soup(["script", "style", "noscript"]):
                        script.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.info(f"Error: {response.status} - {await response.text()}")
    except Exception as e:
        logging.exception(f"Failed to scrape data: {e}")

    prompt = (
        f"Generate a essay for a student in class {data.student_class} with a {data.accent} accent, "
        f"on the topic '{data.topic}', and the mood is '{data.mood}' and give me essay should be less than 400 words "
        f"and in response did not want \n\n or \n and also not want word count thanks you this type of stuff and used {text} "
        f"content for as updated data from internet and which is helpful in created essay and please give me output in paragraph format only not in points."
    )

    username = user.get("username")
    topic = Topic()
    response_text = await topic.topic_data_model_for_Qwen(username, prompt)

    essay_data = FirestoreEssay(
        username=username,
        user_id=user["user_id"],
        student_class=data.student_class,
        accent=data.accent,
        topic=data.topic,
        mood=data.mood,
        content=response_text
    )

    write_time, doc_ref = db.collection("essays").add(essay_data.to_dict())
    essay_id = doc_ref.id

    return JSONResponse(content={
        "response": response_text,
        "essay_id": essay_id
    })


@app.get("/overall-scoring-by-id")
async def overall_scoring_by_id(essay_id: str):
    topic = Topic()
    result = await topic.overall_scoring_by_id(essay_id)
    return result

TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)





@app.websocket("/ws/audio")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]

    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    logging.info(f"[WS] Authenticated connection from {username}")
    chunk_index = 0
    chunk_results = []
    text_output = []

    date_str = datetime.now().strftime("%Y-%m-%d")
    user_dir = os.path.join(TEMP_DIR, username, date_str)
    os.makedirs(user_dir, exist_ok=True)

    final_output = os.path.join(user_dir, f"{username}_output.wav")
    transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

    if os.path.exists(final_output):
        os.remove(final_output)
    if os.path.exists(transcript_path):
        os.remove(transcript_path)

    loop = asyncio.get_event_loop()

    try:
        topic = Topic()
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                print(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
                audio = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                audio.export(chunk_filename, format="wav")

                async with aiohttp.ClientSession() as session:
                    # Emotion detection
                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), content_type="audio/wav")
                        async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
                            emotion_data = await res.json()
                            emotion = emotion_data.get("emotion")

                    # Speech-to-text
                    transcribed_text = await topic.speech_to_text(chunk_filename, username)

                    # Fluency
                    #result = await topic.speech_to_text(audio_file, username)
                    #transcribed_text = result['text']  # Exactly equivalent to old behavior
                    #confidence = result['confidence']
                    async with session.post(
                        f"{CPU_API_BASE}/fluency-score",
                        json={"text": transcribed_text}
                    ) as res:
                        fluency_data = await res.json()
                        fluency = fluency_data.get("fluency")

                    # Pronunciation
                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), content_type="audio/wav")
                        async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
                            pron_data = await res.json()
                            pronunciation = pron_data.get("pronunciation")

                # Silvero remains local
                silvero = await topic.silvero_vad(chunk_filename)
                topic.update_realtime_stats(fluency, pronunciation, emotion)

                text_output.append(transcribed_text)

                chunk_result = {
                    "chunk_index": chunk_index,
                    "text": transcribed_text,
                    "emotion": emotion,
                    "fluency": fluency,
                    "pronunciation": pronunciation,
                    "silvero": silvero,
                    "file_path": chunk_filename
                }

                logging.info(f"[Chunk {chunk_index}] {chunk_result}")
                chunk_results.append(chunk_result)
                chunk_index += 1

    except WebSocketDisconnect:
        logging.warning(f"[WS] {username} forcibly disconnected.")

    finally:
        await loop.run_in_executor(None, merge_chunks, chunk_results, final_output)

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(text_output).strip())

        try:
            essays_ref = db.collection("essays").where("username", "==", username)
            essays = essays_ref.stream()
            today = datetime.now().date()
            latest_essay = max((doc for doc in essays if doc.create_time.date() == today), key=lambda d: d.create_time, default=None)

            if latest_essay:
                essay_ref = db.collection("essays").document(latest_essay.id)
                average_scores = topic.get_average_realtime_scores()
                essay_ref.update({
                    "chunks": chunk_results,
                    "average_scores": average_scores
                })
                logging.info(f"Updated essay {latest_essay.id}")

        except Exception as e:
            logging.error(f"[Firestore Update Error] {e}")

        for file in os.listdir(user_dir):
            if file.startswith("chunk_") and file.endswith(".wav"):
                try:
                    os.remove(os.path.join(user_dir, file))
                except Exception as e:
                    logging.warning(f"Failed to remove {file}: {e}")


def merge_chunks(chunk_files, final_output):
    logging.info("[Merge] Merging audio chunks...")
    combined = AudioSegment.empty()

    for chunk in chunk_files:
        file_path = chunk.get("file_path")
        if file_path and os.path.exists(file_path):
            audio = AudioSegment.from_file(file_path, format="wav")
            combined += audio
        else:
            logging.warning(f"[Merge] Skipping missing or invalid file: {file_path}")

    combined.export(final_output, format="wav")
    logging.info("[Merge] Merged audio file saved.")




@app.get("/get-tts-audio")
def get_tts_audio(username: str):
    folder = os.path.join("text_to_speech_audio_folder", username)
    file_path = os.path.join(folder, f"{username}_output.wav")

    timeout = 60
    poll_interval = 2
    waited = 0

    while waited < timeout:
        if os.path.exists(file_path):
            return FileResponse(file_path, media_type="audio/wav", filename=f"{username}_output.wav")
        time.sleep(poll_interval)
        waited += poll_interval

    raise HTTPException(status_code=408, detail="Audio file not generated within 1 minute.")



genai.configure(api_key=GOOGLE_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-pro")

pc = Pinecone(api_key=PINECONE_API_KEY)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1024,
        metric="cosine",
        pods=1,
        pod_type="p1.x1"
    )
    logging.info(f"Created new Pinecone index {PINECONE_INDEX_NAME} with dimension 1024")

index = pc.Index(PINECONE_INDEX_NAME)

embedding_model = HuggingFaceEmbeddings(
    model_name=EMBEDDING_MODEL_NAME,
    model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'}
)





SUPPORTED_IMAGE_FORMATS = [".png", ".jpg", ".jpeg"]
SUPPORTED_TEXT_FORMATS = [".txt"]
SUPPORTED_PDF_FORMATS = [".pdf"]
SUPPORTED_DOC_FORMATS = [".docx"]
SUPPORTED_PPT_FORMATS = [".pptx"]
SUPPORTED_XLS_FORMATS = [".xlsx"]


def render_text_to_image(text: str, width=800, font_size=18) -> Image.Image:
    font = ImageFont.load_default()
    lines = []
    dummy_img = Image.new("RGB", (width, 1000))
    draw = ImageDraw.Draw(dummy_img)

    words = text.split()
    line = ""
    for word in words:
        test_line = f"{line} {word}".strip()
        bbox = draw.textbbox((0, 0), test_line, font=font)
        w = bbox[2] - bbox[0]
        if w < width - 40:
            line = test_line
        else:
            lines.append(line)
            line = word
    lines.append(line)

    height = font_size * len(lines) + 50
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    y = 20
    for line in lines:
        draw.text((20, y), line, font=font, fill="black")
        y += font_size
    return img


def file_to_images(file_path: str) -> List[Image.Image]:
    ext = os.path.splitext(file_path)[1].lower()
    images = []

    if ext in SUPPORTED_PDF_FORMATS:
        images = convert_from_path(file_path, dpi=200)

    elif ext in SUPPORTED_IMAGE_FORMATS:
        images = [Image.open(file_path)]

    elif ext in SUPPORTED_DOC_FORMATS:
        try:
            text = docx2txt.process(file_path)
            img = render_text_to_image(text)
            images = [img]
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"DOCX to image failed: {e}")

    elif ext in SUPPORTED_PPT_FORMATS:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            img = render_text_to_image(text)
            images.append(img)

    elif ext in SUPPORTED_XLS_FORMATS:
        try:
            excel = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in excel.items():
                text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                img = render_text_to_image(text)
                images.append(img)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"XLSX to image failed: {e}")

    elif ext in SUPPORTED_TEXT_FORMATS:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        img = render_text_to_image(content)
        images = [img]

    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")

    return images


def extract_text_from_any_file(file_path: str) -> str:
    images = file_to_images(file_path)
    all_text = ""

    for idx, image in enumerate(images):
        try:
            response = gemini_model.generate_content([
                "Extract all the text from this image accurately, including tables and special formatting.",
                image
            ])
            text = response.text.strip()
            print(text)
            all_text += f"\n\n--- Page/Image {idx + 1} ---\n{text}"
        except Exception as e:
            logging.error(f"OCR failed on image {idx + 1}: {e}")
            all_text += f"\n\n--- Page/Image {idx + 1} FAILED ---"

    return all_text



@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...),
    student_class: str = Form(...),
    subject: str = Form(...),
    curriculum: str = Form(...)
):
    try:
        # ===== 1. File Storage Setup =====
        folder = f"uploads/{curriculum}/{student_class}/{subject}"
        os.makedirs(folder, exist_ok=True)
        file_path = os.path.join(folder, file.filename)

        # Save file temporarily
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # ===== 2. Text Extraction =====
        extracted_text = extract_text_from_any_file(file_path)
        if not extracted_text.strip():
            raise HTTPException(status_code=422, detail="No text extracted.")

        # ===== 3. Check for Existing File (Deduplication) =====
        namespace = f"{curriculum}_{student_class}_{subject}"
        existing_entries = index.query(
            vector=embedding_model.embed_query(extracted_text[:1000]),
            top_k=1,
            filter={
                "filename": {"$eq": file.filename},
                "type": {"$eq": "ocr_file"}
            },
            namespace=namespace
        )

        if existing_entries.matches:
            logging.warning(f"File {file.filename} already exists in DB")
            return {
                "status": "skipped",
                "message": "File already exists in database",
                "existing_id": existing_entries.matches[0].id
            }

        # ===== 4. Chunking and Metadata Preparation =====
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(extracted_text)

        metadatas = [{
            "curriculum": curriculum,
            "student_class": student_class,
            "subject": subject,
            "filename": file.filename,
            "type": "ocr_file",
            "chunk_id": f"{i}_{uuid.uuid4()}",
            "text": chunk
        } for i, chunk in enumerate(chunks)]

        # ===== 5. Vector Storage with Error Handling =====
        try:
            # Initialize vector store with namespace
            vectorstore = PineconeVectorStore.from_existing_index(
                index_name=PINECONE_INDEX_NAME,
                embedding=embedding_model,
                text_key="text",
                namespace=namespace
            )

            # Generate embeddings and upsert
            embeddings = embedding_model.embed_documents(chunks)
            records = zip(
                [md["chunk_id"] for md in metadatas],
                embeddings,
                metadatas
            )
            
            # Batch upsert to handle large documents
            batch_size = 100  # Pinecone recommends batches of 100-1000 vectors
            for i in range(0, len(chunks), batch_size):
                batch = list(records)[i:i + batch_size]
                vectorstore._index.upsert(vectors=batch)

            # Verify insertion
            stats = index.describe_index_stats()
            new_count = stats["namespaces"].get(namespace, {}).get("vector_count", 0)

        except Exception as e:
            logging.error(f"Pinecone operation failed: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=503,
                detail="Database operation failed. Please retry."
            )

        # ===== 6. Response =====
        return {
            "status": "success",
            "vectors_inserted": new_count,
            "filename": file.filename,
            "namespace": namespace,
            "sample_text": extracted_text[:200] + "..."
        }

    except HTTPException as e:
        raise e

    except Exception as e:
        logging.error(f"Upload failed: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="Internal server error during file processing"
        )

    finally:
        # Clean up temporary file
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                logging.warning(f"Failed to delete temp file: {str(e)}")


def store_in_background(chunks, namespace, filename, curriculum, student_class, subject):
    """Handle storage in background thread"""
    def _store():
        try:
            vectorstore = PineconeVectorStore.from_existing_index(
                index_name=PINECONE_INDEX_NAME,
                embedding=embedding_model,
                text_key="text",
                namespace=namespace
            )

            embeddings = embedding_model.embed_documents(chunks)
            metadatas = [{
                "curriculum": curriculum,
                "student_class": student_class,
                "subject": subject,
                "filename": filename,
                "type": "ocr_file",
                "chunk_id": f"{i}_{uuid.uuid4()}",
                "text": chunk
            } for i, chunk in enumerate(chunks)]

            # Process in small batches
            batch_size = 50
            for i in range(0, len(chunks), batch_size):
                batch = list(zip(
                    [md["chunk_id"] for md in metadatas[i:i+batch_size]],
                    embeddings[i:i+batch_size],
                    metadatas[i:i+batch_size]
                ))
                try:
                    vectorstore._index.upsert(vectors=batch)
                except Exception as e:
                    logging.error(f"Batch upsert failed: {str(e)}")

        except Exception as e:
            logging.error(f"Background storage failed: {str(e)}")

    # Start in background without waiting
    executor.submit(_store)


model_name = OllamaLLM(model="mistral")


async def store_vectors_async(chunks, metadatas, namespace):
    """Background task for storing vectors with error handling"""
    try:
        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )
        
        # Process embeddings in batches
        batch_size = 100
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            
            # Generate embeddings for the batch
            embeddings = await run_in_threadpool(
                embedding_model.embed_documents,
                batch_chunks
            )
            
            # Prepare records for this batch
            batch_metadatas = metadatas[i:i + batch_size]
            records = zip(
                [md["chunk_id"] for md in batch_metadatas],
                embeddings,
                batch_metadatas
            )
            
            # Upsert the batch
            try:
                await run_in_threadpool(
                    vectorstore._index.upsert,
                    vectors=list(records)
                )
                logging.info(f"Upserted batch {i//batch_size + 1}")
            except Exception as e:
                logging.error(f"Failed to upsert batch {i//batch_size + 1}: {str(e)}")
                # Continue with next batch even if one fails
                
    except Exception as e:
        logging.error(f"Vector storage failed: {str(e)}", exc_info=True)

@app.post("/chat/")
async def chat(request: ChatRequest):
    try:
        question = request.question.strip()
        subject = request.subject.strip()
        curriculum = request.curriculum.strip()

        # Create namespace for filtering
        namespace = f"{curriculum}_{subject}"

        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )

        retriever = vectorstore.as_retriever(search_kwargs={
            "k": 5,  # Reduced for better performance
            "filter": {
                "subject": subject,
                "curriculum": curriculum
            }
        })

        prompt_template = PromptTemplate.from_template("""
        You are an expert educator providing clear, concise answers to students.
        Extract the most relevant information to answer the question using ONLY the provided context.

        Follow these rules:
        1. Answer in complete, well-structured sentences.
        2. Do not mention page numbers or document structure.
        3. If context doesn't contain any content, say "This information is not in our materials."
        4. Be factual and avoid speculation.
        5. Use proper grammar and spelling.
        6. Keep your answer concise and to the point.
        7. Do not include '\\n' or '*' in your output.
        8. Do not include escape characters like \\n, \\, \", \n or any slashes.
        9. Do not use markdown symbols like '*', '-', '`', or backslashes.

        Context: {context}
        Question: {question}
        Answer:
        """)

        qa_chain = RetrievalQA.from_chain_type(
            llm=model_name,
            chain_type="stuff",
            retriever=retriever,
            chain_type_kwargs={"prompt": prompt_template},
            return_source_documents=True
        )

        # Corrected: Pass the question directly instead of the chain
        result = qa_chain.invoke({"query": question})  # Using invoke() instead of run()

        return {
            "question": question,
            "answer": result["result"],
            "source_documents": [doc.metadata for doc in result["source_documents"]]
        }

    except Exception as e:
        logging.error(f"Chat error: {str(e)}", exc_info=True)
        return {
            "question": request.question,
            "answer": "An error occurred while processing your request.",
            "error": str(e)
        }


@app.get("/health")
def welcome_page():
    return {"Message": "Welcome the ai speech module page."}




chat_history = []

model = ChatOllama(
            model="mistral", 
            model_kwargs={"temperature": 0.8}
        )

async def system_message(topic, mood, student_class, level) -> SystemMessage:
    parser = StrOutputParser()

    prompt_template = PromptTemplate(template="""
    You are a friendly, knowledgeable teaching assistant. Your purpose is to:
    1. Introduce the topic in a friendly manner
    2. Answer questions conversationally
    3. Never reveal internal project details
    4. Keep responses under 100 words
                                     
    Topic: {topic}
    Mood: {mood}
    Student Class: {student_class}
    Level: {level}
    5. Introduce on understanding the topic always.""",      
    input_variables=["topic", "mood", "student_class", "level"])

    chain = prompt_template | model | parser
    result = await chain.ainvoke({
        "topic": topic,
        "mood": mood,
        "student_class": student_class,
        "level": level
    })
    return result



EMP_DIR = os.path.abspath("temp_chunks")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/assistant")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]
    student_topic = query_params.get("topic", [None])[0]
    student_class = query_params.get("student_class", [None])[0]
    mood = query_params.get("mood", [None])[0]
    accent = query_params.get("accent", [None])[0]

    ai_response = await system_message(student_topic, mood, student_class, accent)

    chat_history = []
    chat_history.append(SystemMessage(content=ai_response))

    
    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    logging.info(f"[WS] Authenticated connection from {username}")
    
    config = {
        'silence_threshold': 2.0,
        'min_utterance_length': 3,
        'max_chunk_duration': 10.0,
        'blocklist': [
            'you', 'thank you', 'tchau', 'thanks', 'ok', 'Obrigado.', 'E aí', '',
            'me', 'hello', 'hi', 'hey', 'okay', 'thanks', 'thank', 'obrigado',
            'tchau.', 'bye', 'goodbye', 'me.', 'you.', 'thank you.',"I'm going to take a picture of the sea","Kansai International Airport",
	    "Thank you for watching!","1 tbsp of salt",'Teksting av Nicolai Winther',

        ],
        'max_repetitions': 2,
        'max_silence': 10.0,
        'chunk_duration': 0.5
    }

    session_state = {
        'silvero_model':True,
        'audio_buffer': AudioSegment.empty(),  # For accumulating speech chunks
        'text_buffer': [],  # For accumulating transcribed text
        'silence_duration': 0.0,
        'last_speech_time': time.time(),
        'conversation_active': False,
        'processing_active': False,  # To prevent overlap
        'chunk_index': 0,
        'chunk_results': [],
        'text_output': []
    }
    topic = Topic()
    session_temp_dir = tempfile.mkdtemp(prefix=f"{username}_", dir=TEMP_DIR)

    response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
    sleep_time = await send_audio_response(websocket, response_audio)
    asyncio.sleep(sleep_time)

    logging.info(f"[Session] Temp dir created at {session_temp_dir}")
    
    final_output = os.path.join(session_temp_dir, f"{username}_output.wav")
    transcript_path = os.path.join(session_temp_dir, f"{username}_transcript.txt")

    try:
        
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                logging.info(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                if session_state['processing_active']:
                    continue
                    
                current_time = time.time()
                
                # Add new audio to buffer
                new_chunk = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                session_state['audio_buffer'] += new_chunk
                
                # Save temp chunk for VAD
                chunk_filename = os.path.join(session_temp_dir, f"chunk_temp_{session_state['chunk_index']}.wav")
                new_chunk.export(chunk_filename, format="wav")
                session_state['chunk_index'] += 1
                
                vad_result = await topic.silvero_vad(chunk_filename)
                current_silence = vad_result.get("duration", 0.0)
                
                if current_silence > 0.3:
                    session_state['silence_duration'] = min(
                        session_state['silence_duration'] + current_silence,
                        config['max_silence']
                    )
                    logging.info(f"[VAD] Silence detected: {current_silence:.2f}s (Total: {session_state['silence_duration']:.2f}s)")
                    
                    if session_state['silence_duration'] >= config['silence_threshold']:

                        await process_buffered_audio(session_state, websocket, username, session_temp_dir, topic, config,student_topic,student_class,mood,accent,chat_history)

    except WebSocketDisconnect:
        logging.warning(f"[WS] {username} disconnected unexpectedly")
    except Exception as e:
        logging.error(f"Unexpected error: {str(e)}", exc_info=True)
    finally:
        essay_id = await finalize_session(session_state, username, session_temp_dir, topic)
        return {"essay_id":essay_id}






def cleanup_temp_files(session_temp_dir, chunk_results):
    """Clean up temporary files from session"""
    try:
        # Remove individual chunk files
        for chunk in chunk_results:
            try:
                if os.path.exists(chunk['file_path']):
                    os.remove(chunk['file_path'])
            except Exception as e:
                logging.warning(f"Failed to delete chunk file {chunk['file_path']}: {e}")

        # Remove the temporary directory if empty
        try:
            if os.path.exists(session_temp_dir):
                if not os.listdir(session_temp_dir):
                    os.rmdir(session_temp_dir)
        except Exception as e:
            logging.warning(f"Failed to remove temp directory {session_temp_dir}: {e}")

    except Exception as e:
        logging.error(f"Error during temp file cleanup: {str(e)}")



async def process_buffered_audio(session_state, websocket, username, temp_dir, topic, config,student_topic,student_class,mood,accent,chat_history):
    if len(session_state['audio_buffer']) == 0:
        return
        
    session_state['processing_active'] = True
    try:
        # Save buffered audio
        buffer_filename = os.path.join(temp_dir, f"buffered_{time.time()}.wav")
        session_state['audio_buffer'].export(buffer_filename, format="wav")
        
        # Transcribe
        transcribed_text = await topic.speech_to_text(buffer_filename, username)
        logging.info(f"Transcribed: {transcribed_text}")
        session_state['text_buffer'].append(transcribed_text)
        session_state['text_output'].append(transcribed_text)
        
        # Save combined audio
        final_output = os.path.join(temp_dir, f"{username}_output.wav")
        if os.path.exists(final_output):
            existing_audio = AudioSegment.from_wav(final_output)
            combined_audio = existing_audio + session_state['audio_buffer']
        else:
            combined_audio = session_state['audio_buffer']
        combined_audio.export(final_output, format="wav")
        
        # Save transcript
        transcript_path = os.path.join(temp_dir, f"{username}_transcript.txt")
        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(session_state['text_output']).strip())
        
        # Process only if meaningful speech
        clean_text = transcribed_text.lower().strip()
        if (len(transcribed_text.split()) >= config['min_utterance_length'] and 
            clean_text not in config['blocklist']):
            
            # Get speech analysis
            async with aiohttp.ClientSession() as session:
                emotion = await detect_emotion(session, buffer_filename)
                fluency = await get_fluency_score(session, transcribed_text)
                pronunciation = await get_pronunciation_score(session, buffer_filename)
                
                # Process the utterance
                await process_user_utterance(
                    transcribed_text, emotion, fluency, pronunciation,
                    session_state, buffer_filename, websocket, 
                    username, temp_dir, topic,student_topic,student_class,mood,accent,chat_history
                )
        
        # Reset buffers
        session_state['audio_buffer'] = AudioSegment.empty()
        session_state['text_buffer'] = []
        session_state['silence_duration'] = 0.0
        session_state['conversation_active'] = True
        
    except Exception as e:
        logging.error(f"Error processing buffered audio: {str(e)}", exc_info=True)
    finally:
        session_state['processing_active'] = False

async def detect_emotion(session, audio_file):
    """Send audio file to emotion detection API"""
    try:
        form = aiohttp.FormData()
        form.add_field('file', 
                      open(audio_file, 'rb'),
                      filename=os.path.basename(audio_file),
                      content_type='audio/wav')
        
        async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as response:
            if response.status == 200:
                data = await response.json()
                return data.get('emotion')
            else:
                logging.error(f"Emotion detection failed: {response.status}")
                return None
    except Exception as e:
        logging.error(f"Emotion detection error: {str(e)}")
        return None
    

async def get_fluency_score(session, text):
    """Send text to fluency scoring API"""
    try:
        async with session.post(
            f"{CPU_API_BASE}/fluency-score",
            json={"text": text}
        ) as response:
            if response.status == 200:
                data = await response.json()
                return data.get('fluency')
            else:
                logging.error(f"Fluency scoring failed: {response.status}")
                return None
    except Exception as e:
        logging.error(f"Fluency scoring error: {str(e)}")
        return None
    


async def scraping(topic: str) -> str:
    url = f"https://en.wikipedia.org/wiki/{topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()  # ✅ FIXED
                    soup = BeautifulSoup(html, "html.parser")
                    for tag in soup(["script", "style", "noscript"]):
                        tag.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.warning(f"[SCRAPING] Failed with status {response.status}: {await response.text()}")
    except Exception as e:
        logging.exception(f"[SCRAPING ERROR] Failed to scrape data: {e}")

    return text




async def get_pronunciation_score(session, audio_file):
    """Send audio file to pronunciation scoring API"""
    try:
        form = aiohttp.FormData()
        form.add_field('file', 
                      open(audio_file, 'rb'),
                      filename=os.path.basename(audio_file),
                      content_type='audio/wav')
        
        async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as response:
            if response.status == 200:
                data = await response.json()
                return data.get('pronunciation')
            else:
                logging.error(f"Pronunciation scoring failed: {response.status}")
                return None
    except Exception as e:
        logging.error(f"Pronunciation scoring error: {str(e)}")
        return None

async def process_user_utterance(text, emotion, fluency, pronunciation, 
                               session_state, chunk_filename, websocket, 
                               username, session_temp_dir, topic,student_topic,student_class,mood,accent,chat_history):
    """Process a valid user utterance"""
    session_state['text_output'].append(text)
    session_state['chunk_results'].append({
        "chunk_index": session_state['chunk_index'],
        "text": text,
        "emotion": emotion,
        "fluency": fluency,
        "pronunciation": pronunciation,
        "file_path": chunk_filename
    })
    session_state['chunk_index'] += 1

    scraping_data = await scraping(student_topic)

    try:
        
        prompt_template = PromptTemplate(template = """
            ROLE: You are a friendly, knowledgeable assistant and communicate like a friend. Your purpose is to:
            1. Answer questions conversationally
            2. Never reveal internal project details
            3. Keep responses under 200 words
            4. Used {scraping_data}, which is updated data related to the topic which help in better answer and talking
            5. if by chance question is much more specific topic is needed and you does not have correct answer then used the function await scraping(your_specific_topic), which help in give updated answer.
                                         
            Chat History:
            {chat_history}
            considering history as well before answering the question.

            Topic: {student_topic}
            Mood: {mood}
            Student Class: {student_class}
            Level: {level}
                                         
            - for reply anything consider always the Topic, Mood , student_class, Level and history.
            - And should be reply in the way and look like a friendly, knowledgeable teaching assistant.
            - Do not mention any technical details, model architecture, or team members.
            - Focus on providing helpful, concise answers.
            - If the question is not related to the topic, politely tell that quesiton is not related to the toic                     
                                                                

            RULES:
            - NEVER mention:
            * Model architecture/type
            * Team members/credentialsS
            * Code implementation
            * Technical specifications
            - Always redirect technical questions to general knowledge                                                                                                 

            USER QUESTION: {question}

            RESPONSE FORMAT:
            [Answer Concise 1-10 sentence response] 
            [if needed  Optional follow-up question to continue conversation]

            EXAMPLE:
            This exaple is for you never asked that. 
            User: What model are you using?
            I focus on helping with learning concepts rather than technical details. 
            Would you like me to explain how these systems generally work?

            Current response should be:
            """,input_variables=["scraping_data","question","student_topic","student_class","mood","accent","chat_history"]
            )
        
        model = ChatOllama(
                model="mistral",
                model_kwargs={"temperature": 0.8}
            )

        parser = StrOutputParser()

        chain = prompt_template | model | parser

        ai_response = await chain.ainvoke({
            "scraping_data":scraping_data,
            "student_topic": student_topic,
            "mood": mood,
            "student_class": student_class,
            "level": accent,
            "question": text,
            "chat_history": chat_history
        })
        chat_history.append(AIMessage(content=ai_response))
        print("[AI Response]:", ai_response)

        response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
        sleep_time = await send_audio_response(websocket, response_audio)
        time.sleep(sleep_time)
        session_state["silvero_model"]=True
        print("session state now : ",session_state["silvero_model"])

    except Exception as e:
        logging.error(f"QA Error: {str(e)}")
        await send_default_response(websocket, username, session_temp_dir, topic)

async def send_audio_response(websocket, audio_file):
    """Send audio file through websocket"""
    try:
        audio = AudioSegment.from_wav(audio_file)
        duration_ms = len(audio)
        with open(audio_file, "rb") as f:
            await websocket.send_bytes(f.read())

        return duration_ms / 1000
    except Exception as e:
        logging.error(f"Failed to send audio response: {str(e)}")

async def send_default_response(websocket, username, session_temp_dir, topic):
    """Send default error response"""
    try:
        default_response = "I didn't quite catch that. Could you please repeat?"
        response_audio = await topic.text_to_speech_assistant(default_response, username, session_temp_dir)
        await send_audio_response(websocket, response_audio)
    except Exception as e:
        logging.error(f"Failed to send default response: {str(e)}")

async def send_followup(websocket, username, session_temp_dir, topic):
    """Send follow-up message after silence"""
    try:
        followup = "Is there anything else I can help you with?"
        response_audio = await topic.text_to_speech_assistant(followup, username, session_temp_dir)
        await send_audio_response(websocket, response_audio)
    except Exception as e:
        logging.error(f"Failed to send follow-up: {str(e)}")

async def finalize_session(session_state, username, session_temp_dir, topic):
    """Cleanup and save session data"""
    essay_id = None
    try:
        essays_ref = db.collection("essays").where("username", "==", username)
        today = datetime.now().date()
        latest_essay = None
        
        docs = essays_ref.stream()
        for doc in docs:
            doc_date = doc.create_time.date() if hasattr(doc, 'create_time') else datetime.fromtimestamp(doc.create_time.seconds).date()
            if doc_date == today:
                if not latest_essay or doc.create_time > latest_essay.create_time:
                    latest_essay = doc

        if latest_essay:
            essay_id = latest_essay.id
            latest_essay.reference.update({
                "chunks": session_state['chunk_results'],
                "average_scores": topic.get_average_realtime_scores(),
                "updated_at": datetime.now()
            })
            logging.info(f"Updated essay document {essay_id}")
            
    except Exception as e:
        logging.error(f"Database update failed: {str(e)}")

    for chunk in session_state['chunk_results']:
        try:
            if os.path.exists(chunk['file_path']):
                os.remove(chunk['file_path'])
        except Exception as e:
            logging.warning(f"Failed to delete chunk file {chunk['file_path']}: {e}")

    return essay_id
